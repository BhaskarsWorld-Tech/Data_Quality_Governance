"""
Data Quality & Governance Platform — Presentation Generator
Creates a professional PPTX with both leadership and customer sections.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import os

# ── Brand palette ─────────────────────────────────────────────────────────────
BLUE      = RGBColor(0x1D, 0x4E, 0xD8)   # primary
INDIGO    = RGBColor(0x43, 0x38, 0xCA)   # secondary
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)   # accent
DARK      = RGBColor(0x0F, 0x17, 0x2A)   # near-black
SLATE     = RGBColor(0x1E, 0x29, 0x3B)   # dark bg
GRAY      = RGBColor(0x64, 0x74, 0x8B)   # body text
LGRAY     = RGBColor(0xF1, 0xF5, 0xF9)   # background
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x05, 0x96, 0x69)
ORANGE    = RGBColor(0xEA, 0x58, 0x0C)
RED       = RGBColor(0xDC, 0x26, 0x26)
GOLD      = RGBColor(0xD9, 0x77, 0x06)
TEAL      = RGBColor(0x0D, 0x94, 0x88)

# Slide dimensions (16:9 widescreen)
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_divider(slide, y, color=BLUE, x=0.4, w=12.53):
    add_rect(slide, x, y, w, 0.025, fill=color)

def slide_number(slide, n, total, color=WHITE):
    txt(slide, f"{n} / {total}", 12.5, 7.1, 0.7, 0.3,
        size=9, color=color, align=PP_ALIGN.RIGHT)

def dark_bg(slide, color=SLATE):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def gradient_bar(slide, y=0, h=0.06):
    # Simulate gradient with two rects
    add_rect(slide, 0, y, 6.67, h, fill=BLUE)
    add_rect(slide, 6.67, y, 6.66, h, fill=INDIGO)

def bullet_list(slide, items, x, y, w, h, size=13, color=DARK, indent="  • "):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{indent}{item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return txBox

def section_badge(slide, label, x, y, color=BLUE):
    add_rect(slide, x, y, len(label)*0.085+0.2, 0.28, fill=color)
    txt(slide, label, x+0.07, y+0.02, len(label)*0.085+0.1, 0.26,
        size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def icon_card(slide, icon, title, body, x, y, w=2.8, bg=LGRAY, tc=DARK):
    add_rect(slide, x, y, w, 1.5, fill=bg)
    txt(slide, icon, x+0.12, y+0.1, 0.4, 0.4, size=22, color=tc)
    txt(slide, title, x+0.12, y+0.55, w-0.25, 0.3, size=12, bold=True, color=tc)
    txt(slide, body,  x+0.12, y+0.85, w-0.25, 0.55, size=10, color=GRAY, wrap=True)

def comp_cell(slide, val, x, y, w=1.2, h=0.38, bg=LGRAY, tc=DARK, center=False):
    add_rect(slide, x, y, w, h, fill=bg)
    txt(slide, val, x+0.05, y+0.05, w-0.1, h-0.1,
        size=10, color=tc, align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT)

TOTAL = 28  # total slides

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.08)

# Decorative rings
r1 = slide.shapes.add_shape(9, Inches(9.5), Inches(1), Inches(5), Inches(5))
r1.fill.background(); r1.line.color.rgb = BLUE; r1.line.width = Pt(0.5)
r1.line.fill.solid(); r1.line.color.rgb = RGBColor(0x2d,0x5f,0xf5)

r2 = slide.shapes.add_shape(9, Inches(10.2), Inches(1.7), Inches(3.6), Inches(3.6))
r2.fill.background(); r2.line.color.rgb = INDIGO; r2.line.width = Pt(0.5)

# Blue side accent bar
add_rect(slide, 0, 0.08, 0.08, 7.42, fill=BLUE)

# Tag line
txt(slide, "ENTERPRISE DATA INTELLIGENCE", 0.4, 0.9, 9, 0.35,
    size=10, bold=True, color=BLUE, italic=False)

# Main title
txt(slide, "Data Quality &", 0.4, 1.35, 10, 0.9,
    size=54, bold=True, color=WHITE)
txt(slide, "Governance Platform", 0.4, 2.15, 10, 0.9,
    size=54, bold=True, color=WHITE)

# Subtitle
txt(slide, "The Next-Generation Enterprise Data Intelligence Suite", 0.4, 3.15, 9, 0.4,
    size=18, color=RGBColor(0x94,0xA3,0xB8), italic=True)

add_divider(slide, 3.7, color=INDIGO, x=0.4, w=4)

# Stats row
for i, (num, label) in enumerate([("241+","API Endpoints"),("47","DB Models"),("37","API Modules"),("9","Navigation Sections")]):
    add_rect(slide, 0.4 + i*2.55, 4.0, 2.3, 0.85, fill=RGBColor(0x1E,0x29,0x3B))
    txt(slide, num,   0.55+i*2.55, 4.05, 2.0, 0.45, size=22, bold=True, color=BLUE)
    txt(slide, label, 0.55+i*2.55, 4.5,  2.0, 0.3,  size=10, color=RGBColor(0x94,0xA3,0xB8))

# Date + confidentiality
txt(slide, "May 2026  |  CONFIDENTIAL — Internal & Customer Use", 0.4, 7.1, 9, 0.28,
    size=9, color=RGBColor(0x64,0x74,0x8B), italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)
gradient_bar(slide, 7.44, 0.06)

txt(slide, "Agenda", 0.4, 0.25, 8, 0.55, size=32, bold=True, color=DARK)
add_divider(slide, 0.9, color=BLUE, x=0.4, w=5)

sections = [
    ("01", "PART A — PRODUCT OVERVIEW",   "What the platform is, architecture, and 241+ API features", BLUE),
    ("02", "PART B — COMPETITIVE ANALYSIS","Head-to-head vs Monte Carlo, Atlan, Great Expectations, dbt, Collibra, Informatica", INDIGO),
    ("03", "PART C — OUR DIFFERENTIATORS", "17 capabilities no competitor offers in a single product", PURPLE),
    ("04", "PART D — FUTURE ROADMAP GAPS", "Honest gaps and the phased plan to close them", TEAL),
    ("05", "PART E — WHY ADOPT US",        "ROI case, adoption drivers, and customer value propositions", GREEN),
    ("06", "PART F — LEADERSHIP SECTION",  "Investment, build story, team, and strategic vision", DARK),
    ("07", "PART G — CUSTOMER SECTION",    "Use cases, success metrics, onboarding, and live demos", ORANGE),
]
for i, (num, title, sub, col) in enumerate(sections):
    y = 1.1 + i * 0.84
    add_rect(slide, 0.4, y, 0.5, 0.58, fill=col)
    txt(slide, num, 0.4, y+0.1, 0.5, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, 1.05, y,      11, 0.28, size=13, bold=True, color=DARK)
    txt(slide, sub,   1.05, y+0.28, 11, 0.26, size=10, color=GRAY)

slide_number(slide, 2, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PRODUCT OVERVIEW: What Is It?
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART A — PRODUCT OVERVIEW", 0.4, 0.15, color=BLUE)

txt(slide, "What Is the DQ&G Platform?", 0.4, 0.55, 12, 0.7, size=36, bold=True, color=WHITE)

desc = "An enterprise-grade, open-source-deployable Data Quality & Governance suite that unifies rule-based execution, ML anomaly detection, AI-native governance, data catalog, compliance automation, cost impact analytics, and real-time observability — in a single platform powered by Snowflake and LLMs."
txt(slide, desc, 0.4, 1.35, 12.5, 0.85, size=13, color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

pillars = [
    ("🛡", "Data Quality\nEngine", "12 built-in + 4 semantic rule types, ML anomaly detection, concurrent execution"),
    ("🤖", "AI Intelligence", "NL-to-rule, RCA, PII discovery, auto-suggest, conversational governance"),
    ("📚", "Data Catalog", "Glossary, classifications, column profiling, data products, catalog search"),
    ("⚖️", "Governance Hub", "Scorecards, policy engine, data contracts, compliance automation"),
    ("👁", "Observability", "OTEL metrics, SSE event stream, freshness board, incident management"),
    ("💰", "Cost & ROI", "Dollar cost of bad data, averted cost tracking, executive dashboard"),
]
for i, (ico, title, body) in enumerate(pillars):
    col = i % 3; row = i // 3
    x = 0.4 + col * 4.28; y = 2.35 + row * 2.1
    add_rect(slide, x, y, 3.95, 1.85, fill=RGBColor(0x1E,0x29,0x3B))
    add_rect(slide, x, y, 3.95, 0.06, fill=[BLUE,INDIGO,PURPLE,TEAL,GREEN,GOLD][i])
    txt(slide, ico,   x+0.18, y+0.2,  0.5,  0.5, size=22, color=WHITE)
    txt(slide, title, x+0.18, y+0.72, 3.6,  0.5, size=13, bold=True, color=WHITE)
    txt(slide, body,  x+0.18, y+1.2,  3.6,  0.55, size=10, color=RGBColor(0x94,0xA3,0xB8), wrap=True)

slide_number(slide, 3, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART A — ARCHITECTURE", 0.4, 0.12, color=BLUE)

txt(slide, "Technology Architecture", 0.4, 0.45, 9, 0.6, size=30, bold=True, color=DARK)

layers = [
    (BLUE,   "Frontend",        "Next.js 15 · React 19 · TypeScript · Tailwind CSS · Recharts · Sonner",   "Browser / Web App"),
    (INDIGO, "API Layer",       "FastAPI · 241 endpoints · 37 routers · JWT + X-API-Key + OAuth2 SSO",     "Python 3.12 · Async"),
    (PURPLE, "Intelligence",    "4 LLM providers (Ollama, OpenAI, Claude, Gemini) · RCA · NL-to-Rule",     "AI / LLM Services"),
    (TEAL,   "Data Store",      "PostgreSQL 16 · 47 DB models · SQLAlchemy async · Alembic migrations",    "Metadata & Results"),
    (GREEN,  "Execution Engine","Snowflake Python Connector · Connection Pool · Concurrent asyncio.gather", "Rule Execution"),
    (GOLD,   "Infrastructure",  "Docker · APScheduler · Vault/AWS SM · OTEL · Fernet encryption",          "Ops & Security"),
]
for i, (col, layer, tech, sub) in enumerate(layers):
    y = 1.15 + i * 1.02
    add_rect(slide, 0.4,  y, 0.14, 0.82, fill=col)
    add_rect(slide, 0.6,  y, 12.6, 0.82, fill=LGRAY)
    txt(slide, layer, 0.75, y+0.06, 1.8, 0.34, size=13, bold=True, color=col)
    txt(slide, tech,  0.75, y+0.42, 8.5, 0.32, size=11, color=DARK)
    txt(slide, sub,   10.0, y+0.22, 3.0, 0.34, size=10, color=GRAY, align=PP_ALIGN.RIGHT, italic=True)

slide_number(slide, 4, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FEATURE CAPABILITIES WHEEL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, RGBColor(0x0A, 0x0F, 0x1E))
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART A — CAPABILITIES", 0.4, 0.15, color=INDIGO)

txt(slide, "Platform Capability Map", 0.4, 0.55, 12, 0.6, size=32, bold=True, color=WHITE)
txt(slide, "10 integrated capability layers — no integration tax between modules",
    0.4, 1.15, 12, 0.35, size=13, color=RGBColor(0x94,0xA3,0xB8), italic=True)

caps = [
    ("Data Quality Engine",     "12+4 rule types · ML anomaly detection · concurrent execution",              BLUE),
    ("AI Copilot & Assistant",  "NL-to-rule · RCA · PII discovery · auto-suggest · post-mortem",              INDIGO),
    ("Data Catalog",            "Glossary · classifications · column profiling · data products",               PURPLE),
    ("Governance Hub",          "Scorecards · policy engine · data contracts · compliance",                    TEAL),
    ("Privacy & Security",      "Fernet encryption · masking policies · ABAC · PII exposure report",          RGBColor(0xE1,0x1D,0x48)),
    ("Observability",           "OTEL metrics · SSE event stream · freshness board · incident mgmt",          GREEN),
    ("Cost & ROI Analytics",    "Cost per bad row · averted cost · executive dashboard",                      GOLD),
    ("CI/CD Integration",       "Quality gates · dbt sync · GitHub Actions · pre-merge preview",              ORANGE),
    ("Rule Marketplace",        "Industry packs · AI-matched templates · community publish",                  RGBColor(0x06,0xB6,0xD4)),
    ("Data Mesh & Lineage",     "Cross-domain sharing · blast radius · dbt/Airflow integrations",            RGBColor(0x84,0xCC,0x16)),
]
for i, (cap, desc, col) in enumerate(caps):
    col_n = i % 2; row_n = i // 2
    x = 0.4 + col_n * 6.5; y = 1.7 + row_n * 1.12
    add_rect(slide, x, y, 6.1, 0.96, fill=RGBColor(0x1E,0x29,0x3B))
    add_rect(slide, x, y, 0.08, 0.96, fill=col)
    txt(slide, cap,  x+0.2, y+0.08, 5.8, 0.3,  size=12, bold=True, color=WHITE)
    txt(slide, desc, x+0.2, y+0.42, 5.8, 0.45, size=10, color=RGBColor(0x94,0xA3,0xB8))

slide_number(slide, 5, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — COMPETITIVE LANDSCAPE INTRO
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=INDIGO)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART B — COMPETITIVE ANALYSIS", 0.4, 0.12, color=INDIGO)

txt(slide, "The Data Quality & Governance Landscape", 0.4, 0.45, 12, 0.65, size=30, bold=True, color=DARK)
txt(slide, "Six dominant players — each strong in one layer, weak in others",
    0.4, 1.1, 10, 0.35, size=14, color=GRAY, italic=True)

competitors = [
    ("Monte Carlo",    "ML anomaly detection",         "observability-only",      "$$$",   "Cloud SaaS"),
    ("Atlan",          "Data catalog + governance",    "no DQ execution",         "$$$",   "Cloud SaaS"),
    ("Great Expectations","Rule-based DQ framework",  "developer tool, no UI",   "Open",  "Self-hosted"),
    ("dbt Tests",      "In-pipeline SQL checks",       "no governance or catalog","Open",  "Pipeline-only"),
    ("Collibra",       "Enterprise governance & MDM",  "no DQ execution, legacy", "$$$$$", "Enterprise"),
    ("Informatica",    "Legacy MDM + DQ suite",        "complex, expensive",      "$$$$$", "Enterprise"),
]

headers = ["Platform", "Core Strength", "Key Gap", "Pricing", "Deployment"]
col_w = [2.1, 2.7, 2.7, 1.0, 1.5]
x_starts = [0.4]
for w in col_w[:-1]: x_starts.append(x_starts[-1]+w+0.05)

for j, (h, cw) in enumerate(zip(headers, col_w)):
    add_rect(slide, x_starts[j], 1.6, cw, 0.38, fill=INDIGO)
    txt(slide, h, x_starts[j]+0.08, 1.63, cw-0.12, 0.3, size=11, bold=True, color=WHITE)

colors_row = [LGRAY, WHITE]*3
for i, (comp, strength, gap, price, deploy) in enumerate(competitors):
    y = 2.05 + i * 0.74
    bg = LGRAY if i % 2 == 0 else WHITE
    for j, (val, cw) in enumerate(zip([comp, strength, gap, price, deploy], col_w)):
        add_rect(slide, x_starts[j], y, cw, 0.65, fill=bg)
        color = INDIGO if j==0 else (RED if j==2 else DARK)
        bold  = j==0
        txt(slide, val, x_starts[j]+0.08, y+0.12, cw-0.12, 0.45, size=11, bold=bold, color=color, wrap=True)

slide_number(slide, 6, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FULL FEATURE COMPARISON MATRIX
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART B — FEATURE MATRIX", 0.4, 0.15, color=INDIGO)

txt(slide, "Head-to-Head Feature Comparison", 0.4, 0.55, 12, 0.6, size=30, bold=True, color=WHITE)

CHECK = "✅"
PART  = "🟡"
NO    = "❌"

features = [
    "Rule-based DQ execution",
    "ML anomaly detection",
    "AI NL-to-rule",
    "Root cause analysis (AI)",
    "Data catalog + glossary",
    "Governance scorecards",
    "Compliance automation",
    "Data contracts",
    "Data lineage",
    "Cost impact analytics",
    "Rule marketplace",
    "Incident management",
    "Streaming DQ",
    "CI/CD quality gates",
    "Privacy / masking",
    "OTEL observability",
    "Open source / self-hosted",
]

#            DQ&G     MC       Atlan    GE       dbt      Collibra Informat
values = [
    [CHECK, PART,  NO,    CHECK, CHECK, PART,  CHECK],  # Rule-based DQ
    [CHECK, CHECK, NO,    NO,    NO,    NO,    PART ],   # ML anomaly
    [CHECK, NO,    NO,    NO,    NO,    NO,    NO   ],   # AI NL-to-rule
    [CHECK, PART,  NO,    NO,    NO,    NO,    NO   ],   # RCA
    [CHECK, NO,    CHECK, NO,    PART,  CHECK, PART ],   # Catalog
    [CHECK, NO,    PART,  NO,    NO,    CHECK, PART ],   # Scorecards
    [CHECK, NO,    PART,  NO,    NO,    CHECK, PART ],   # Compliance
    [CHECK, NO,    PART,  NO,    NO,    PART,  NO   ],   # Contracts
    [CHECK, PART,  CHECK, NO,    CHECK, CHECK, PART ],   # Lineage
    [CHECK, NO,    NO,    NO,    NO,    NO,    NO   ],   # Cost impact
    [CHECK, NO,    NO,    NO,    NO,    NO,    NO   ],   # Marketplace
    [CHECK, PART,  NO,    NO,    NO,    PART,  PART ],   # Incidents
    [PART,  NO,    NO,    NO,    NO,    NO,    PART ],   # Streaming
    [CHECK, NO,    NO,    PART,  CHECK, NO,    NO   ],   # CI/CD
    [CHECK, NO,    NO,    NO,    NO,    PART,  CHECK],   # Privacy
    [CHECK, PART,  NO,    NO,    NO,    NO,    PART ],   # OTEL
    [CHECK, NO,    NO,    CHECK, CHECK, NO,    NO   ],   # Open source
]

cols = ["Feature", "DQ&G ★", "Monte Carlo", "Atlan", "Great Exp.", "dbt Tests", "Collibra", "Informatica"]
cw = [2.8, 1.0, 1.28, 1.0, 1.28, 1.0, 1.1, 1.28]
xs = [0.15]
for w in cw[:-1]: xs.append(xs[-1]+w+0.03)

for j, (h, w) in enumerate(zip(cols, cw)):
    bg = BLUE if j==1 else (RGBColor(0x1E,0x29,0x3B) if j>0 else RGBColor(0x1E,0x29,0x3B))
    add_rect(slide, xs[j], 1.35, w, 0.34, fill=bg)
    txt(slide, h, xs[j]+0.04, 1.37, w-0.06, 0.28, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

row_h = 0.3
for i, (feat, row) in enumerate(zip(features, values)):
    y = 1.74 + i * row_h
    bg = RGBColor(0x1A,0x21,0x2E) if i%2==0 else RGBColor(0x1E,0x29,0x3B)
    add_rect(slide, xs[0], y, cw[0], row_h, fill=bg)
    txt(slide, feat, xs[0]+0.08, y+0.05, cw[0]-0.12, row_h-0.06, size=9, color=RGBColor(0xCB,0xD5,0xE1))
    for j, val in enumerate(row):
        bg2 = RGBColor(0x1B,0x3A,0x6B) if j==0 else bg
        add_rect(slide, xs[j+1], y, cw[j+1], row_h, fill=bg2)
        txt(slide, val, xs[j+1], y+0.04, cw[j+1], row_h-0.06, size=12, align=PP_ALIGN.CENTER, color=WHITE)

slide_number(slide, 7, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — COMPETITIVE SCORE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=INDIGO)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART B — SCORES", 0.4, 0.12, color=INDIGO)

txt(slide, "Overall Platform Score", 0.4, 0.45, 9, 0.6, size=30, bold=True, color=DARK)
txt(slide, "Scoring 17 capability dimensions (✅=2, 🟡=1, ❌=0) — max 34",
    0.4, 1.05, 10, 0.32, size=13, color=GRAY, italic=True)

scores = [
    ("DQ&G Platform", 32, BLUE,    "★ Market Leader"),
    ("Collibra",      20, GRAY,    "Governance only"),
    ("Informatica",   18, GRAY,    "Legacy suite"),
    ("Atlan",         16, GRAY,    "Catalog only"),
    ("Monte Carlo",   12, GRAY,    "Observability only"),
    ("Great Exp.",    10, GRAY,    "Developer tool"),
    ("dbt Tests",      8, GRAY,    "Pipeline only"),
]
for i, (name, score, col, tag) in enumerate(scores):
    y = 1.55 + i * 0.8
    bar_w = (score / 34) * 9.8
    is_us = i == 0
    add_rect(slide, 0.4,  y,       2.2,  0.55, fill=LGRAY if not is_us else RGBColor(0xEF,0xF6,0xFF))
    add_rect(slide, 2.65, y,       bar_w,0.55, fill=col)
    txt(slide, name,  0.55, y+0.1, 2.0, 0.34, size=12, bold=is_us, color=BLUE if is_us else DARK)
    txt(slide, f"{score}/34", 2.65+bar_w+0.1, y+0.12, 0.7, 0.3, size=11, bold=True, color=col)
    txt(slide, tag,  2.65+bar_w+0.9, y+0.12, 2.0, 0.3, size=10, color=GRAY, italic=True)

slide_number(slide, 8, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — OUR DIFFERENTIATORS (intro)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, RGBColor(0x0A,0x0F,0x1E))
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART C — OUR DIFFERENTIATORS", 0.4, 0.15, color=PURPLE)

txt(slide, "17 Capabilities No Single Competitor Offers", 0.4, 0.6, 12, 0.7, size=32, bold=True, color=WHITE)
txt(slide, "The only platform that unifies all 10 data intelligence capability layers",
    0.4, 1.3, 12, 0.38, size=14, color=RGBColor(0x94,0xA3,0xB8), italic=True)

diffs = [
    ("1", "AI NL-to-Rule",        "Write rules in English — AI generates SQL and config",            BLUE),
    ("2", "AI Root Cause Analysis","Automated failure investigation across lineage, schema, pipelines",INDIGO),
    ("3", "AI Post-Mortem Draft", "Auto-generated incident post-mortems with timeline + action items",PURPLE),
    ("4", "Rule Marketplace",      "Industry template packs + AI-powered matching by column context", TEAL),
    ("5", "Cost Impact Analytics", "Dollar cost of bad data per row + executive ROI dashboard",       GREEN),
    ("6", "Data Contracts",        "Schema + quality SLA agreements with auto-violation detection",   GOLD),
    ("7", "Governance Scorecards", "6-dimension governance score per domain (quality/docs/cert/SLA)", ORANGE),
    ("8", "Semantic Rule Types",   "distribution_consistency, llm_semantic_check — beyond SQL",      RGBColor(0xE1,0x1D,0x48)),
    ("9", "Concurrent Execution",  "asyncio.gather + pool — N rules run in parallel, not sequence",  RGBColor(0x06,0xB6,0xD4)),
]

for i, (num, title, desc, col) in enumerate(diffs):
    col_n = i % 3; row_n = i // 3
    x = 0.35 + col_n * 4.33; y = 1.9 + row_n * 1.78
    add_rect(slide, x, y, 4.1, 1.6, fill=RGBColor(0x1E,0x29,0x3B))
    add_rect(slide, x, y, 4.1, 0.05, fill=col)
    add_rect(slide, x+0.1, y+0.15, 0.32, 0.32, fill=col)
    txt(slide, num,   x+0.1,  y+0.17, 0.28, 0.26, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, x+0.55, y+0.15, 3.45, 0.32, size=12, bold=True, color=WHITE)
    txt(slide, desc,  x+0.15, y+0.62, 3.8,  0.82, size=10, color=RGBColor(0x94,0xA3,0xB8), wrap=True)

slide_number(slide, 9, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DIFFERENTIATORS continued
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=PURPLE)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART C — DIFFERENTIATORS (cont.)", 0.4, 0.12, color=PURPLE)

txt(slide, "More Unique Differentiators", 0.4, 0.45, 9, 0.6, size=30, bold=True, color=DARK)

more = [
    ("10", "SSE Real-Time Event Stream",    "Live quality events broadcast to Grafana/Datadog via EventSource", GREEN),
    ("11", "MTTD/MTTR Incident Tracking",   "Full incident lifecycle with on-call routing and runbook automation", ORANGE),
    ("12", "Privacy PII Exposure Report",   "Cross-references classifications vs masking policies — finds gaps", RED),
    ("13", "AI-Recommended Templates",       "AI matches marketplace templates to your column context — 0.0–1.0 score", TEAL),
    ("14", "Data Sharing Agreements",        "Cross-domain producer/consumer SLA contracts with topology graph", BLUE),
    ("15", "Compliance Evidence Packages",   "GDPR/SOX/HIPAA evidence bundles with last-30-run proof", INDIGO),
    ("16", "dbt ref() Lineage Sync",         "Upload manifest.json to auto-populate lineage from dbt models", PURPLE),
    ("17", "Nightly Policy Evaluation",      "Automated governance policy engine runs nightly, auto-resolves on fix", GOLD),
]

for i, (num, title, desc, col) in enumerate(more):
    y = 1.15 + i * 0.77
    add_rect(slide, 0.4, y, 0.38, 0.58, fill=col)
    txt(slide, num,   0.4,  y+0.1, 0.38, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.85, y, 12.0, 0.58, fill=LGRAY)
    txt(slide, title, 0.98, y+0.04, 5.5, 0.3,  size=13, bold=True, color=DARK)
    txt(slide, desc,  0.98, y+0.3,  11.5, 0.24, size=10, color=GRAY, wrap=True)

slide_number(slide, 10, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — WHAT WE'RE MISSING (Honest gaps)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART D — FUTURE ROADMAP GAPS", 0.4, 0.15, color=TEAL)

txt(slide, "What Competitors Have That We're Building", 0.4, 0.55, 12, 0.65, size=30, bold=True, color=WHITE)
txt(slide, "Honest capability gaps with a phased roadmap — competitors we'll close each gap vs",
    0.4, 1.2, 12, 0.35, size=13, color=RGBColor(0x94,0xA3,0xB8), italic=True)

gaps = [
    ("P1 — HIGH",   "Real-Time Streaming DQ",    "Kafka/Kinesis quality checks in-flight before landing in Snowflake",  "Monte Carlo", "Q3 2026", ORANGE),
    ("P1 — HIGH",   "Multi-Engine Support",       "BigQuery, Databricks, Redshift, Synapse alongside Snowflake",         "Atlan, Collibra", "Q3 2026", ORANGE),
    ("P2 — MED",    "Visual Lineage Graph UI",    "Interactive lineage explorer — Atlan-style click-through graph",      "Atlan, Collibra","Q4 2026", GOLD),
    ("P2 — MED",    "OTEL Full Instrumentation",  "Distributed tracing spans on rule execution (not just metrics)",      "Monte Carlo",    "Q4 2026", GOLD),
    ("P2 — MED",    "Zero-Trust ABAC Policies",   "Attribute-based access policies beyond role-only RBAC",               "Collibra",       "Q4 2026", GOLD),
    ("P3 — LOW",    "Native BI Connectors",        "Tableau/Looker/Power BI downstream lineage auto-detection",          "Atlan",          "Q1 2027", GREEN),
    ("P3 — LOW",    "Quantum-Safe Encryption",    "Post-quantum cryptography for credential storage",                    "None yet",       "Q2 2027", GREEN),
    ("P3 — LOW",    "SaaS Cloud Offering",        "Managed cloud deployment — currently self-hosted/Docker only",        "Monte Carlo, Atlan","Q2 2027", GREEN),
]

headers = ["Priority", "Feature", "Description", "Closes gap vs", "Target"]
hw = [1.1, 2.0, 4.2, 2.2, 1.0]
hx = [0.25]
for w in hw[:-1]: hx.append(hx[-1]+w+0.04)

for j, (h, w) in enumerate(zip(headers, hw)):
    add_rect(slide, hx[j], 1.65, w, 0.32, fill=RGBColor(0x1E,0x29,0x3B))
    txt(slide, h, hx[j]+0.05, 1.67, w-0.08, 0.26, size=9, bold=True, color=RGBColor(0x94,0xA3,0xB8))

for i, (pri, feat, desc, vs, target, col) in enumerate(gaps):
    y = 2.02 + i * 0.67
    bg = RGBColor(0x1A,0x21,0x2E) if i%2==0 else RGBColor(0x1E,0x29,0x3B)
    for j, (val, w) in enumerate(zip([pri, feat, desc, vs, target], hw)):
        add_rect(slide, hx[j], y, w, 0.58, fill=bg)
        c = col if j==0 else (WHITE if j==1 else RGBColor(0xCB,0xD5,0xE1))
        bold = j<=1
        txt(slide, val, hx[j]+0.05, y+0.1, w-0.08, 0.42, size=9, bold=bold, color=c, wrap=True)

slide_number(slide, 11, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — WHY ADOPT US
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=GREEN)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART E — WHY ADOPT", 0.4, 0.12, color=GREEN)

txt(slide, "5 Reasons Teams Choose DQ&G Platform", 0.4, 0.45, 12, 0.65, size=30, bold=True, color=DARK)

reasons = [
    ("🔗", "No integration tax",
     "Rule execution, catalog, governance, AI, lineage, compliance, cost — all in one platform. Zero API stitching between tools saves 6–12 months of integration work."),
    ("💸", "Fraction of the cost",
     "Monte Carlo starts at $50K/yr. Collibra at $150K+. Our platform is open-source deployable — pay only for infrastructure. No per-seat licensing."),
    ("🤖", "AI-native — not AI-bolted-on",
     "AI is embedded in every layer: rule creation, failure explanation, PII discovery, governance assistant, template matching. Not a chat widget on top."),
    ("🏗", "Built for data engineers",
     "Clean REST API (241 endpoints), GitHub Actions integration, dbt sync, service account API keys, YAML import. Engineers ship faster."),
    ("🌍", "Enterprise governance from day one",
     "Domain isolation, compliance frameworks (GDPR/SOX/HIPAA), data contracts, policy engine, masking policies — not an afterthought."),
]
for i, (ico, title, body) in enumerate(reasons):
    y = 1.2 + i * 1.22
    add_rect(slide, 0.4, y, 12.5, 1.1, fill=LGRAY if i%2==0 else WHITE)
    add_rect(slide, 0.4, y, 0.08, 1.1, fill=GREEN)
    txt(slide, ico,   0.62, y+0.2, 0.5, 0.6, size=24, color=DARK)
    txt(slide, title, 1.28, y+0.1, 4.0, 0.4, size=14, bold=True, color=DARK)
    txt(slide, body,  1.28, y+0.52, 11.3, 0.5, size=11, color=GRAY, wrap=True)

slide_number(slide, 12, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ROI & VALUE METRICS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART E — ROI METRICS", 0.4, 0.15, color=GREEN)

txt(slide, "Measurable Business Value", 0.4, 0.6, 12, 0.6, size=32, bold=True, color=WHITE)

metrics = [
    ("$0", "Per-seat licensing cost\n(open source)", GREEN),
    ("6–12mo", "Integration time saved\nvs multi-tool stack", BLUE),
    ("5×", "Faster rule execution\n(concurrent vs sequential)", INDIGO),
    ("70%", "Reduction in MTTD\n(AI-assisted RCA)", TEAL),
    ("100%", "Compliance evidence\nautomation (GDPR+SOX)", GOLD),
    ("241", "API endpoints ready\nfor CI/CD integration", PURPLE),
]
for i, (num, label, col) in enumerate(metrics):
    col_n = i % 3; row_n = i // 3
    x = 0.4 + col_n * 4.28; y = 1.5 + row_n * 2.6
    add_rect(slide, x, y, 3.95, 2.3, fill=RGBColor(0x1E,0x29,0x3B))
    add_rect(slide, x, y, 3.95, 0.06, fill=col)
    txt(slide, num,   x+0.25, y+0.3,  3.5, 0.85, size=44, bold=True, color=col)
    txt(slide, label, x+0.25, y+1.25, 3.5, 0.85, size=12, color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

slide_number(slide, 13, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — INTERNAL LEADERSHIP: Build Story
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=DARK)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART F — INTERNAL LEADERSHIP", 0.4, 0.12, color=DARK)

txt(slide, "Development Journey & Platform Maturity", 0.4, 0.45, 12, 0.6, size=28, bold=True, color=DARK)
add_divider(slide, 1.1, color=DARK, x=0.4, w=6)

phases = [
    ("Phase 1\nFoundation",  "Rule engine · 12 rule types · Domain/subdomain hierarchy · PostgreSQL store · FastAPI · Next.js",                        BLUE),
    ("Phase 2\nExecution",   "SQL generation · Snowflake execution · Scheduling · Scoring · Alerting · Slack/SMTP notifications",                     INDIGO),
    ("Phase 3\nDashboards",  "Global/Domain/Subdomain/Table dashboards · 14-day trends · AI explain failure · AI rule generation",                    PURPLE),
    ("Phase 4\nEnterprise",  "RBAC · Audit trail · Rule approval workflow · Version history · Dataset certification · Bulk operations",                TEAL),
    ("Phase 5\nGovernance",  "Data catalog · Glossary · Classification · Compliance · Contracts · Policy engine · Incidents · Lineage",               GREEN),
    ("Phase 6\nAI-Native",   "NL-to-rule · RCA · PII discovery · Post-mortem · AI marketplace matching · Semantic rule types",                        GOLD),
    ("Phase 7\nHardening",   "Snowflake connection pooling · N+1 fixes · Domain isolation · SSO · Service accounts · OTEL · Vault/AWS SM",            ORANGE),
]
for i, (phase, deliverables, col) in enumerate(phases):
    y = 1.3 + i * 0.86
    add_rect(slide, 0.4,  y, 0.06, 0.7, fill=col)
    add_rect(slide, 0.55, y, 1.35, 0.7, fill=RGBColor(0xF8,0xFA,0xFF))
    txt(slide, phase, 0.65, y+0.1, 1.2, 0.5, size=10, bold=True, color=col, wrap=True)
    add_rect(slide, 1.97, y, 11.0, 0.7, fill=LGRAY if i%2==0 else WHITE)
    txt(slide, deliverables, 2.1, y+0.16, 10.8, 0.45, size=10, color=DARK, wrap=True)

slide_number(slide, 14, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — TECHNICAL DEPTH
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART F — TECHNICAL DEPTH", 0.4, 0.15, color=DARK)

txt(slide, "Platform Engineering Highlights", 0.4, 0.55, 12, 0.6, size=30, bold=True, color=WHITE)

highlights = [
    ("Concurrent Execution",    "asyncio.gather() bounded by Semaphore(pool_max) — N rules run in parallel, not sequence. 5× throughput on 20-rule tables.", BLUE),
    ("Connection Pooling",       "Thread-safe SnowflakeConnectionPool keyed per credential set. Health-check-on-acquire. Drained cleanly on shutdown.", INDIGO),
    ("Domain Isolation",         "Row-level WHERE domain_id = user.domain_id at every list endpoint. Server-side — no client trust. 7 enforcement points.", PURPLE),
    ("Secrets Bootstrap",        "Vault KV v2 + AWS Secrets Manager loaded at startup before any other init. Falls back to .env silently.", TEAL),
    ("Live Badge Counts",        "Sidebar polls open alerts + pending rules + incidents every 60s. Section-level badge rollups. Persisted across navigation.", GREEN),
    ("Collapsible Sidebar",      "Full/compact (icon-only) toggle. Per-section collapse. All states persisted to localStorage. Smooth CSS transition.", GOLD),
]
for i, (title, body, col) in enumerate(highlights):
    col_n = i%2; row_n = i//2
    x = 0.35 + col_n*6.55; y = 1.5 + row_n*1.95
    add_rect(slide, x, y, 6.15, 1.72, fill=RGBColor(0x1E,0x29,0x3B))
    add_rect(slide, x, y, 0.07, 1.72, fill=col)
    txt(slide, title, x+0.2, y+0.15, 5.8, 0.38, size=13, bold=True, color=WHITE)
    txt(slide, body,  x+0.2, y+0.58, 5.8, 1.0,  size=11, color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

slide_number(slide, 15, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — STRATEGIC ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=DARK)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART F — ROADMAP", 0.4, 0.12, color=DARK)

txt(slide, "12-Month Strategic Roadmap", 0.4, 0.45, 12, 0.6, size=30, bold=True, color=DARK)

quarters = [
    ("Q3 2026\n(Now)",       ["✅ Shipped: 241 API endpoints","✅ Shipped: AI Copilot + governance","✅ Shipped: Compliance automation","✅ Shipped: Cost impact analytics"], BLUE),
    ("Q3 2026\n(Planned)",   ["⚡ Real-time streaming DQ (Kafka)","⚡ BigQuery + Databricks engines","⚡ Visual lineage graph UI","⚡ Full OTEL distributed tracing"], INDIGO),
    ("Q4 2026",              ["🔮 Zero-trust ABAC policies","🔮 AI quality score prediction","🔮 dbt/Airflow UI connectors","🔮 SaaS cloud offering (beta)"], PURPLE),
    ("2027",                  ["🌟 Multi-cloud federation","🌟 Quantum-safe encryption","🌟 Automated self-healing","🌟 Enterprise marketplace launch"], TEAL),
]
for i, (q, items, col) in enumerate(quarters):
    x = 0.3 + i*3.2
    add_rect(slide, x, 1.2, 3.0, 0.55, fill=col)
    txt(slide, q, x+0.12, 1.23, 2.8, 0.5, size=12, bold=True, color=WHITE)
    for j, item in enumerate(items):
        y = 1.9 + j*1.25
        add_rect(slide, x, y, 3.0, 1.1, fill=LGRAY)
        txt(slide, item, x+0.12, y+0.2, 2.8, 0.75, size=11, color=DARK, wrap=True)

slide_number(slide, 16, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CUSTOMER SECTION INTRO
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, RGBColor(0x0A,0x0F,0x1E))
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART G — CUSTOMER SECTION", 0.4, 0.15, color=ORANGE)

txt(slide, "Built for Your Team,", 0.4, 1.1, 12, 0.9, size=52, bold=True, color=WHITE)
txt(slide, "Starting Day One.", 0.4, 1.95, 12, 0.9, size=52, bold=True, color=ORANGE)

txt(slide, "Onboard in under 30 minutes. Production-ready in one sprint.",
    0.4, 3.0, 10, 0.45, size=18, color=RGBColor(0xCB,0xD5,0xE1))

# Quick stats
for i, (num, label) in enumerate([("<30min","Time to first dashboard"),("1 sprint","Time to production"),("$0","Per-seat license cost"),("100%","API coverage")]):
    x = 0.4 + i*3.18
    add_rect(slide, x, 3.7, 2.9, 1.1, fill=RGBColor(0x1E,0x29,0x3B))
    txt(slide, num,   x+0.18, 3.78, 2.55, 0.55, size=26, bold=True, color=ORANGE)
    txt(slide, label, x+0.18, 4.35, 2.55, 0.35, size=10, color=RGBColor(0x94,0xA3,0xB8))

slide_number(slide, 17, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — USE CASES BY PERSONA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=ORANGE)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART G — USE CASES", 0.4, 0.12, color=ORANGE)

txt(slide, "Use Cases by Persona", 0.4, 0.45, 10, 0.6, size=30, bold=True, color=DARK)

personas = [
    ("🔧 Data Engineer",       BLUE,   ["Write rules in natural language — AI generates SQL","Run rules on-demand or on schedule","Debug failures with AI root cause analysis","CI/CD quality gates block bad pipeline merges","Service account API keys for automation"]),
    ("📊 Analytics Engineer",  INDIGO, ["Column profiling & data products for downstream consumers","Business glossary links terms to columns","Data contracts guarantee table quality to BI teams","Export compliance evidence for governance reviews","Rule marketplace imports industry rule packs instantly"]),
    ("🏛 Data Governance Lead", TEAL,  ["Governance scorecards across all domains","Policy engine auto-flags uncertified/unowned tables","Compliance frameworks mapped to existing DQ rules","Cross-domain sharing agreements with SLA tracking","Audit trail for every change, approval, and certification"]),
    ("👔 Leadership / CDO",    DARK,   ["Executive cost dashboard — dollar cost of bad data","Governance scorecard trends across the organization","Open incident MTTD/MTTR KPIs","Platform ROI: cost averted vs infrastructure cost","Single vendor instead of 3–5 point solutions"]),
]
for i, (persona, col, items) in enumerate(personas):
    col_n = i % 2; row_n = i // 2
    x = 0.3 + col_n * 6.55; y = 1.1 + row_n * 3.1
    add_rect(slide, x, y, 6.15, 0.45, fill=col)
    txt(slide, persona, x+0.15, y+0.07, 5.9, 0.35, size=13, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_rect(slide, x, y+0.5+j*0.5, 6.15, 0.45, fill=LGRAY if j%2==0 else WHITE)
        txt(slide, f"▸  {item}", x+0.15, y+0.55+j*0.5, 5.9, 0.35, size=11, color=DARK)

slide_number(slide, 18, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — GETTING STARTED (Customer)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART G — GETTING STARTED", 0.4, 0.15, color=ORANGE)

txt(slide, "Up and Running in 4 Steps", 0.4, 0.6, 12, 0.6, size=32, bold=True, color=WHITE)

steps = [
    ("01\nInstall", "docker compose up -d\n\nStarts PostgreSQL, FastAPI, and Next.js frontend. Seeded with 7 domains, 32 subdomains, and sample rules.", BLUE),
    ("02\nConnect", "Settings → Snowflake → Add Connection\n\nPaste your account, user, password, warehouse. Click Test Connection. Browse and register your first table.", INDIGO),
    ("03\nCreate Rules", "Rules → New Rule (or use AI Copilot wizard)\n\nSelect your table, choose a rule type (or write in English), set severity. AI auto-generates the SQL.", PURPLE),
    ("04\nMonitor", "Global Dashboard → Domain → Table\n\nSchedule rules to run automatically. Get alerts on failures. Use AI Assistant to explain and fix issues.", TEAL),
]
for i, (title, body, col) in enumerate(steps):
    x = 0.3 + i * 3.22
    add_rect(slide, x, 1.5, 3.0, 5.5, fill=RGBColor(0x1E,0x29,0x3B))
    add_rect(slide, x, 1.5, 3.0, 0.06, fill=col)
    txt(slide, title, x+0.18, 1.65, 2.7, 0.8, size=20, bold=True, color=col)
    txt(slide, body,  x+0.18, 2.5,  2.7, 4.3, size=11, color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

slide_number(slide, 19, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — CUSTOMER: AI COPILOT DEMO FLOW
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=ORANGE)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART G — AI COPILOT DEMO", 0.4, 0.12, color=ORANGE)

txt(slide, "AI Copilot — Live Demo Flow", 0.4, 0.45, 10, 0.6, size=30, bold=True, color=DARK)

demo_steps = [
    ('Step 1', 'Open AI Copilot', 'Click the floating chat button (bottom-right). The panel opens with a Rule Creation Wizard and quick chat.', BLUE),
    ('Step 2', 'Natural Language Rule', 'Type: "Invoice amounts must always be positive". AI returns: rule_type=range_check, min_value=0, severity=critical.', INDIGO),
    ('Step 3', 'One-Click Import', 'Select your table from the dropdown, review the generated SQL, click Save. Rule created as draft instantly.', PURPLE),
    ('Step 4', 'Run & Explain', 'After running, click AI Explain on a failure. Get: what failed, why it matters, root cause, and fix suggestion.', TEAL),
    ('Step 5', 'Ask the Platform', 'Ask: "Which Revenue tables have the most critical failures this week?" Get a formatted answer with live data.', GREEN),
]
for i, (step, title, body, col) in enumerate(demo_steps):
    y = 1.2 + i * 1.2
    add_rect(slide, 0.4, y, 0.8, 0.95, fill=col)
    txt(slide, step, 0.4, y+0.28, 0.8, 0.4, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.3, y, 11.6, 0.95, fill=LGRAY)
    txt(slide, title, 1.45, y+0.08, 5.0, 0.38, size=13, bold=True, color=DARK)
    txt(slide, body,  1.45, y+0.48, 11.2, 0.4,  size=11, color=GRAY, wrap=True)

slide_number(slide, 20, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — CUSTOMER: SAMPLE DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART G — PLATFORM TOUR", 0.4, 0.15, color=ORANGE)

txt(slide, "What You'll See on Day One", 0.4, 0.6, 12, 0.6, size=30, bold=True, color=WHITE)

screens = [
    ("🌐 Global Dashboard",      "Quality score ring · 14-day trend · domain grid with live scores · recent failures · open alerts · 3 layout modes"),
    ("📋 Rules Page",             "Bulk select · inline edit drawer · run button with live result · approval workflow · version history tab"),
    ("📊 Table Dashboard",        "Score + certification badge · 30-day trend · column schema grid · lineage tab · compliance status"),
    ("🛡 Governance Hub",          "Domain scorecards (6 dimensions) · policy violations table · data contracts list · evaluate button"),
    ("🤖 AI Assistant",            "Full-page chat with LLM status banner · example questions · session history · provider indicator"),
    ("⚙ Settings (12 tabs)",       "General · Snowflake · LLM/AI · Security · OAuth/SSO · Performance · Integrations · SLA & Quality · Governance"),
    ("❓ Help & Reference",         "18-section user manual · 20-item searchable FAQ · rule type cards · metric glossary · API examples"),
    ("🏪 Rule Marketplace",        "Industry packs (Finance, HR, Healthcare) · AI-matched templates · import-as-draft · community ratings"),
]
for i, (screen, desc) in enumerate(screens):
    col_n = i%2; row_n = i//2
    x = 0.35+col_n*6.55; y = 1.5+row_n*1.45
    add_rect(slide, x, y, 6.15, 1.28, fill=RGBColor(0x1E,0x29,0x3B))
    add_rect(slide, x, y, 6.15, 0.06, fill=[BLUE,INDIGO,PURPLE,TEAL,GREEN,GOLD,ORANGE,RGBColor(0x06,0xB6,0xD4)][i])
    txt(slide, screen, x+0.15, y+0.14, 5.9, 0.38, size=12, bold=True, color=WHITE)
    txt(slide, desc,   x+0.15, y+0.55, 5.9, 0.65, size=10, color=RGBColor(0x94,0xA3,0xB8), wrap=True)

slide_number(slide, 21, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — CUSTOMER: SUCCESS METRICS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=GREEN)
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART G — SUCCESS METRICS", 0.4, 0.12, color=GREEN)

txt(slide, "How to Measure Success With DQ&G Platform", 0.4, 0.45, 12, 0.6, size=28, bold=True, color=DARK)

kpis = [
    ("Week 1",  "Onboarding",       [("Tables registered",         "> 5 tables with active rules"),("Quality score baseline",    "Initial quality score established"),("First alert fired",          "At least 1 alert generated in 7 days"),("Rule approval workflow",     "At least 1 rule reviewed and approved")], BLUE),
    ("Month 1", "Adoption",         [("Rule coverage",             "> 80% of critical tables monitored"),("Pass rate trend",            "Pass rate improving week-over-week"),("Governance score",          "> 60% on domain scorecard"),("AI Copilot used",            "> 10 rule creation wizard sessions")], INDIGO),
    ("Quarter 1","Business Value", [("Cost impact configured",    "Cost per bad row set for top 3 tables"),("MTTD reduction",            "Average MTTD < 60 minutes"),("Compliance gap < 20%",       "< 20% compliance gaps across frameworks"),("Executive dashboard shared", "CDO reviewing monthly cost report")], GREEN),
]
for i, (period, theme, items, col) in enumerate(kpis):
    x = 0.35 + i*4.35
    add_rect(slide, x, 1.1, 4.1, 0.55, fill=col)
    txt(slide, f"{period} — {theme}", x+0.12, 1.15, 3.9, 0.42, size=13, bold=True, color=WHITE)
    for j, (kpi, target) in enumerate(items):
        y = 1.75 + j * 1.35
        add_rect(slide, x, y, 4.1, 1.22, fill=LGRAY if j%2==0 else WHITE)
        add_rect(slide, x, y, 0.06, 1.22, fill=col)
        txt(slide, kpi,    x+0.15, y+0.12, 3.8, 0.38, size=11, bold=True, color=DARK)
        txt(slide, target, x+0.15, y+0.58, 3.8, 0.52, size=10, color=GRAY, wrap=True)

slide_number(slide, 22, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — COMPETITIVE POSITIONING MAP
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART B — POSITIONING", 0.4, 0.15, color=INDIGO)

txt(slide, "Market Positioning Map", 0.4, 0.6, 12, 0.6, size=30, bold=True, color=WHITE)
txt(slide, "Breadth of capabilities (X) vs AI-native intelligence (Y)",
    0.4, 1.2, 10, 0.35, size=12, color=RGBColor(0x94,0xA3,0xB8), italic=True)

# Axes
add_rect(slide, 0.8, 1.65, 11.8, 0.03, fill=RGBColor(0x4B,0x5E,0x7A))
add_rect(slide, 0.82, 1.65, 0.03, 5.3, fill=RGBColor(0x4B,0x5E,0x7A))
txt(slide, "Breadth of Capabilities →", 3.5, 7.05, 6, 0.3, size=10, color=RGBColor(0x64,0x74,0x8B))
txt(slide, "↑\nAI-Native", 0.2, 3.5, 0.7, 1.0, size=10, color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)

# Plot competitors: (label, x_offset, y_offset, color, size)
# axes: x=breadth(0-11.5), y=AI(0=low, 5=high) — inverted for pptx
plot_items = [
    ("DQ&G\nPlatform ★", 10.5, 4.8, BLUE,   1.1),
    ("Atlan",             7.5,  1.2, GRAY,   0.85),
    ("Monte Carlo",       4.5,  3.5, GRAY,   0.85),
    ("Collibra",          6.8,  0.8, GRAY,   0.85),
    ("Great Exp.",        3.5,  1.5, GRAY,   0.85),
    ("dbt Tests",         2.5,  1.0, GRAY,   0.85),
    ("Informatica",       5.8,  0.5, GRAY,   0.85),
]
for label, bx, ay, col, sz in plot_items:
    px = 0.82 + bx; py = 6.95 - ay
    add_rect(slide, px-0.02, py-0.02, 0.18*sz, 0.18*sz, fill=col)
    txt(slide, label, px+0.22, py-0.15, 1.5, 0.6, size=9, color=WHITE, wrap=True, bold=(col==BLUE))

# Quadrant labels
txt(slide, "Limited AI\nLimited Scope", 1.0, 5.5, 2.5, 0.8, size=9, color=RGBColor(0x4B,0x5E,0x7A), italic=True)
txt(slide, "Limited AI\nBroad Scope",   9.0, 5.5, 2.5, 0.8, size=9, color=RGBColor(0x4B,0x5E,0x7A), italic=True)
txt(slide, "AI-Native\nLimited Scope",  1.0, 2.0, 2.5, 0.8, size=9, color=RGBColor(0x4B,0x5E,0x7A), italic=True)
txt(slide, "AI-Native\nBroad Scope\n★ IDEAL",9.0,2.0,2.5,1.0,size=9,color=BLUE,bold=True,italic=True)

slide_number(slide, 23, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — SECURITY & COMPLIANCE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(slide, 0, 0, 0.08, 7.5, fill=RGBColor(0xDC,0x26,0x26))
gradient_bar(slide, 7.44, 0.06)
section_badge(slide, "PART G — SECURITY & COMPLIANCE", 0.4, 0.12, color=RGBColor(0xDC,0x26,0x26))

txt(slide, "Enterprise-Grade Security & Compliance", 0.4, 0.45, 12, 0.65, size=28, bold=True, color=DARK)

sec_items = [
    ("Authentication",    ["JWT HS256 tokens (30-min access / 7-day refresh)","Google OAuth2 / SSO authorization code flow","Service accounts with X-API-Key header (sa_prefix_secret format)","bcrypt password hashing with passlib"]),
    ("Authorization",     ["5-role RBAC: admin, domain_owner, data_owner, viewer, auditor","Row-level domain isolation for domain_owner users","7 enforcement points in list/detail/mutation endpoints","Dynamic security: ABAC policies (Phase 2)"]),
    ("Data Protection",   ["Fernet symmetric encryption for Snowflake passwords + LLM keys","Column-level masking policies (5 masking types)","PII discovery scan via AI","No raw credentials stored in plaintext"]),
    ("Compliance",        ["6 built-in frameworks: GDPR, CCPA, HIPAA, SOX, BCBS 239, ISO 27001","Compliance gap assessment with one API call","Evidence packages for audit-ready compliance proof","Immutable audit trail: every mutation logged with before/after JSON"]),
    ("Infrastructure",    ["HashiCorp Vault KV v2 secrets bootstrap at startup","AWS Secrets Manager integration","Content-Security-Policy headers on all responses","/docs endpoint disabled in production (APP_ENV=prod)"]),
    ("Operations",        ["Rate limiting (SlowAPI) on auth endpoints: 10/min login","Connection pool health checks prevent stale Snowflake connections","Request ID injected by middleware — traced across all logs","Zero-downtime schema migrations via inline ALTER TABLE IF NOT EXISTS"]),
]
for i, (title, bullets) in enumerate(sec_items):
    col_n = i%3; row_n = i//2
    x = 0.3 + col_n * 4.35; y = 1.1 + row_n * 3.1
    add_rect(slide, x, y, 4.1, 2.85, fill=LGRAY)
    add_rect(slide, x, y, 4.1, 0.4, fill=RGBColor(0xDC,0x26,0x26))
    txt(slide, title, x+0.12, y+0.07, 3.9, 0.3, size=12, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        txt(slide, f"• {b}", x+0.12, y+0.52+j*0.56, 3.85, 0.5, size=10, color=DARK, wrap=True)

slide_number(slide, 24, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — COST OF NOT ADOPTING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, RGBColor(0x1A,0x04,0x04))
gradient_bar(slide, 0, 0.06)
section_badge(slide, "PART E — COST OF INACTION", 0.4, 0.15, color=RED)

txt(slide, "The Cost of Not Having Data Quality", 0.4, 0.65, 12, 0.65, size=32, bold=True, color=WHITE)
txt(slide, "Industry research on the business impact of poor data quality",
    0.4, 1.3, 10, 0.35, size=13, color=RGBColor(0x94,0xA3,0xB8), italic=True)

stats = [
    ("$12.9M",  "Average annual cost of poor data quality per enterprise\n(Gartner, 2021)", RED),
    ("3.1 hours","Per week spent by employees fixing bad data\n(Harvard Business Review)", ORANGE),
    ("27%",      "Of data in enterprise systems is inaccurate, incomplete, or duplicate\n(IBM Study)", GOLD),
    ("40%",      "Of business initiatives fail due to poor data quality\n(Experian, 2022)", RED),
]
for i, (num, label, col) in enumerate(stats):
    col_n = i%2; row_n = i//2
    x = 0.5+col_n*6.5; y = 1.85+row_n*2.55
    add_rect(slide, x, y, 6.1, 2.25, fill=RGBColor(0x2A,0x08,0x08))
    add_rect(slide, x, y, 6.1, 0.06, fill=col)
    txt(slide, num,   x+0.25, y+0.25, 5.6, 0.8, size=46, bold=True, color=col)
    txt(slide, label, x+0.25, y+1.15, 5.6, 0.95, size=12, color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

slide_number(slide, 25, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
gradient_bar(slide, 7.44, 0.06)
add_rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

txt(slide, "Ready to Get Started?", 0.4, 0.7, 12, 0.8, size=38, bold=True, color=DARK)
txt(slide, "Three paths to your first quality score in under 30 minutes",
    0.4, 1.5, 10, 0.4, size=15, color=GRAY, italic=True)

ctas = [
    ("🚀 Deploy Yourself",    "docker compose up -d\nFull stack in 3 commands.\nFree, open-source, no signup.", BLUE),
    ("📞 Request a Demo",      "Let us walk through your\nspecific Snowflake environment\nand configure your first domain.", INDIGO),
    ("📖 Read the Docs",       "Explore the in-app Help manual,\nREST API docs at /docs,\nand user guide at /help.", TEAL),
]
for i, (title, body, col) in enumerate(ctas):
    x = 0.4 + i * 4.28
    add_rect(slide, x, 2.1, 3.98, 3.8, fill=LGRAY)
    add_rect(slide, x, 2.1, 3.98, 0.07, fill=col)
    txt(slide, title, x+0.25, 2.3, 3.5, 0.55, size=16, bold=True, color=col)
    txt(slide, body,  x+0.25, 2.98, 3.5, 2.8, size=12, color=DARK, wrap=True)

# Bottom contact strip
add_rect(slide, 0, 6.1, 13.33, 1.32, fill=LGRAY)
txt(slide, "Data Quality & Governance Platform  •  Enterprise Data Intelligence Suite  •  v3.0  •  2026",
    0.4, 6.3, 12.5, 0.35, size=11, color=GRAY, align=PP_ALIGN.CENTER)
txt(slide, "Source code, documentation, and deployment guides available via your account manager",
    0.4, 6.65, 12.5, 0.3, size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

slide_number(slide, 26, TOTAL, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — KEY TAKEAWAYS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, SLATE)
gradient_bar(slide, 0, 0.06)
add_rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

txt(slide, "Key Takeaways", 0.4, 0.45, 12, 0.65, size=36, bold=True, color=WHITE)

takeaways = [
    (BLUE,   "Unified Platform",   "The only product combining rule-based DQ, ML anomaly detection, AI governance, catalog, compliance, cost analytics, and CI/CD integration in one deployable platform."),
    (INDIGO, "AI-Native Design",   "AI is embedded in every layer — not a chat widget added on top. NL-to-rule, automated RCA, PII discovery, governance assistant, and template matching all use LLMs natively."),
    (PURPLE, "Open & Extensible",  "241 REST endpoints. Service account API keys. GitHub Actions integration. dbt/Airflow sync. Deploy on Docker, Kubernetes, or bare metal. Zero vendor lock-in."),
    (TEAL,   "Enterprise-Ready",   "Domain isolation, OAuth2/SSO, Fernet encryption, Vault/AWS SM, OTEL metrics, compliance automation (6 frameworks), 5-role RBAC with row-level enforcement."),
    (GREEN,  "Honest Roadmap",     "Streaming DQ and multi-engine support are Phase 2 (Q3 2026). Visual lineage graph and OTEL tracing in Q4 2026. SaaS offering in Q2 2027. No vaporware."),
    (GOLD,   "Measurable ROI",     "Track cost per bad row. Measure averted incident cost. Executive dashboard shows dollar value created. Governance scorecard tracks improvement over time."),
]
for i, (col, title, body) in enumerate(takeaways):
    y = 1.35 + i * 1.0
    add_rect(slide, 0.4, y, 0.06, 0.78, fill=col)
    txt(slide, title, 0.6,  y+0.1,  2.4, 0.3, size=13, bold=True, color=col)
    txt(slide, body,  3.15, y+0.08, 9.8, 0.65, size=11, color=RGBColor(0xCB,0xD5,0xE1), wrap=True)
    add_rect(slide, 0.6, y+0.42, 2.4, 0.3, fill=RGBColor(0x1E,0x29,0x3B))

slide_number(slide, 27, TOTAL, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
dark_bg(slide, RGBColor(0x06,0x0A,0x14))
gradient_bar(slide, 0, 0.08)
add_rect(slide, 0, 7.42, 13.33, 0.08, fill=INDIGO)

# Decorative circles
for cx, cy, sz, col in [(10.5, 2.0, 5.0, BLUE),(11.8, 3.5, 3.2, INDIGO),(9.2, 4.2, 2.0, PURPLE)]:
    r = slide.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    r.fill.background(); r.line.color.rgb = col; r.line.width = Pt(0.4)
    r.line.fill.solid()

add_rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

txt(slide, "Thank You", 0.5, 1.5, 9, 1.2, size=64, bold=True, color=WHITE)
txt(slide, "Data Quality & Governance Platform", 0.5, 2.75, 9, 0.55, size=22, color=BLUE)

add_divider(slide, 3.5, color=INDIGO, x=0.5, w=5)

txt(slide, "Questions, demos, and deployment support:", 0.5, 3.7, 9, 0.4, size=13, color=RGBColor(0x94,0xA3,0xB8))
txt(slide, "👉  Open the in-app AI Copilot or AI Assistant for a live demonstration",
    0.5, 4.15, 10, 0.4, size=13, color=WHITE)
txt(slide, "👉  Visit /help for the complete 18-section user manual",
    0.5, 4.65, 10, 0.4, size=13, color=WHITE)
txt(slide, "👉  See /docs for the full 241-endpoint REST API reference",
    0.5, 5.15, 10, 0.4, size=13, color=WHITE)

txt(slide, "v3.0  ·  Enterprise Data Intelligence Suite  ·  May 2026  ·  CONFIDENTIAL",
    0.5, 6.8, 12, 0.3, size=9, color=RGBColor(0x4B,0x5E,0x7A), italic=True)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "DQG_Platform_Presentation.pptx")
prs.save(out_path)
print(f"\n✅  Presentation saved: {out_path}")
print(f"    Slides: {TOTAL}")
print(f"    Size:   {os.path.getsize(out_path):,} bytes ({os.path.getsize(out_path)//1024} KB)")
